import subprocess
import tempfile
import os
import time

# Step 1: Create and open PowerPoint
code_str = """
from pptx import Presentation
import os

prs = Presentation()
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Yoga: A Journey to Wellness"
slide.placeholders[1].text = "An Introduction to Yoga Practices and Benefits"

prs.save("Yoga2Presentation.pptx")
os.system("open -a 'Microsoft PowerPoint' Yoga2Presentation.pptx")
"""

with tempfile.NamedTemporaryFile(suffix=".py", delete=False, mode='w') as f:
    f.write(code_str)
    temp_filename = f.name

subprocess.run(["python3", temp_filename])
time.sleep(3)  # wait for PowerPoint to open

# Step 2: AppleScript to export as JPEG
apple_script = '''
tell application "Microsoft PowerPoint"
    activate
    delay 1
    set thePres to active presentation
    set imgFolder to (POSIX path of (path to desktop folder)) & "slide_images"
    export thePres to imgFolder as PNG
end tell
'''
subprocess.run(["osascript", "-e", apple_script])

# Step 3: Open the exported JPEG
jpeg_path = os.path.expanduser("~/Desktop/slide_images/Slide1.JPG")
if os.path.exists(jpeg_path):
    subprocess.run(["open", jpeg_path])
else:
    print("Slide1.JPG not found.")
