from pptx import Presentation
from pptx.util import Inches
import os
import uuid
from glob import glob

SAVE_DIR = "saved_pptx"
os.makedirs(SAVE_DIR, exist_ok=True)

def save_pptx(slides, pptx_id):
    prs = Presentation()
    for slide in slides:
        s = prs.slides.add_slide(prs.slide_layouts[5])
        title = s.shapes.title
        title.text = slide.get("title", "")
        body = s.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(4))
        body.text_frame.text = slide.get("content", "")
    
    path = os.path.join(SAVE_DIR, f"{pptx_id}.pptx")
    prs.save(path)
    return path

def assemble_all_pptx():
    final = Presentation()
    files = glob(os.path.join(SAVE_DIR, "*.pptx"))
    
    for file in files:
        src = Presentation(file)
        for slide in src.slides:
            slide_copy = final.slides.add_slide(final.slide_layouts[5])
            for shape in slide.shapes:
                if shape.has_text_frame:
                    textbox = slide_copy.shapes.add_textbox(
                        Inches(1), Inches(1), Inches(6), Inches(4)
                    )
                    textbox.text_frame.text = shape.text
    
    output_path = os.path.join(SAVE_DIR, "final_presentation.pptx")
    final.save(output_path)
    return output_path
