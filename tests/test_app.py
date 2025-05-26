import unittest
from unittest.mock import patch, MagicMock
from app import app
import os

class TestFlaskApp(unittest.TestCase):

    def setUp(self):
        self.app = app.test_client()
        self.app.testing = True

    def test_index_route(self):
        response = self.app.get("/")
        self.assertEqual(response.status_code, 200)
        self.assertIn(b"<html", response.data)

    @patch("app.client.chat.completions.create")
    def test_generate_pptx_success(self, mock_openai):
        # Mock OpenAI returning working pptx code
        mock_response = MagicMock()
        mock_response.choices = [MagicMock(message=MagicMock(content="""
            from pptx import Presentation
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            title = slide.shapes.title
            title.text = "Test Slide"
            prs.save("output.pptx")
                    """))]
        mock_openai.return_value = mock_response

        response = self.app.post("/generate-pptx", data={"user_prompt": "Make a test slide"})

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.content_type, "application/vnd.openxmlformats-officedocument.presentationml.presentation")

        # Clean up generated file
        if os.path.exists("output.pptx"):
            os.remove("output.pptx")

    @patch("app.client.chat.completions.create")
    def test_generate_pptx_openai_failure(self, mock_openai):
        mock_openai.side_effect = Exception("OpenAI is down")
        response = self.app.post("/generate-pptx", data={"user_prompt": "Will this fail?"})

        self.assertEqual(response.status_code, 500)
        self.assertIn(b"OpenAI is down", response.data)

if __name__ == '__main__':
    unittest.main()
